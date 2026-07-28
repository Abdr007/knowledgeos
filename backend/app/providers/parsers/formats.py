"""Concrete parsers for PDF, DOCX, PPTX, Markdown/text and HTML.

One module: each parser is a dozen lines of library glue, and five files whose
entire content is "call this library and map the result" is filing, not
architecture. The protocol boundary is what matters and it is preserved —
adding OCR still means adding one class and one registry entry (§29.1).
"""

from __future__ import annotations

import io
import logging
import re

from app.core.errors import UnsupportedMediaError, ValidationError
from app.db.models.enums import SourceType
from app.providers.parsers.base import ParsedDocument, ParsedPage

logger = logging.getLogger(__name__)

_WS = re.compile(r"[ \t\x0b\f\r]+")
_BLANKS = re.compile(r"\n{3,}")


def normalize(text: str) -> str:
    """Collapse whitespace without destroying paragraph structure.

    Paragraph breaks survive because the chunker splits on them (§11); flattening
    all whitespace here would remove the strongest boundary signal the chunker has.
    """
    text = text.replace("\u00a0", " ").replace("\r\n", "\n").replace("\r", "\n")
    text = _WS.sub(" ", text)
    text = "\n".join(line.strip() for line in text.split("\n"))
    return _BLANKS.sub("\n\n", text).strip()


class PdfParser:
    source_type = SourceType.PDF

    def parse(self, data: bytes, *, filename: str | None = None) -> ParsedDocument:
        from pypdf import PdfReader
        from pypdf.errors import PdfReadError

        try:
            reader = PdfReader(io.BytesIO(data))
        except PdfReadError as exc:
            raise ValidationError(f"Could not read the PDF: {exc}") from exc

        if reader.is_encrypted:
            # Attempt the empty password, which is what "encrypted but not
            # password-protected" means in practice for many exported PDFs.
            try:
                if reader.decrypt("") == 0:
                    raise ValidationError("The PDF is password-protected.")
            except (NotImplementedError, ValidationError):
                raise ValidationError("The PDF is password-protected.") from None

        pages: list[ParsedPage] = []
        for index, page in enumerate(reader.pages, start=1):
            try:
                text = normalize(page.extract_text() or "")
            except Exception:
                logger.warning("failed to extract a PDF page", extra={"page": index})
                text = ""
            if text:
                pages.append(ParsedPage(number=index, text=text))

        meta: dict[str, object] = dict(reader.metadata or {})
        title = str(meta.get("/Title") or "").strip() or None
        return ParsedDocument(
            pages=pages, title=title, metadata={"pages": str(len(reader.pages))}
        )


class DocxParser:
    source_type = SourceType.DOCX

    def parse(self, data: bytes, *, filename: str | None = None) -> ParsedDocument:
        import docx

        try:
            document = docx.Document(io.BytesIO(data))
        except Exception as exc:
            raise ValidationError(f"Could not read the DOCX file: {exc}") from exc

        # DOCX has no pages until it is rendered, so paragraphs are grouped under
        # their nearest heading and each heading becomes one "page". That gives
        # citations a meaningful anchor ("Methodology") instead of a fake number.
        sections: list[tuple[str | None, list[str]]] = [(None, [])]
        for para in document.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            if para.style is not None and str(para.style.name).startswith("Heading"):
                sections.append((text, []))
            else:
                sections[-1][1].append(text)

        for table in document.tables:
            rows = [
                " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                for row in table.rows
            ]
            rows = [r for r in rows if r]
            if rows:
                sections[-1][1].append("\n".join(rows))

        pages = [
            ParsedPage(number=i, text=normalize("\n\n".join(body)), section=heading)
            for i, (heading, body) in enumerate(sections, start=1)
            if body
        ]
        core = document.core_properties
        return ParsedDocument(pages=pages, title=(core.title or "").strip() or None)


class PptxParser:
    source_type = SourceType.PPTX

    def parse(self, data: bytes, *, filename: str | None = None) -> ParsedDocument:
        from pptx import Presentation

        try:
            deck = Presentation(io.BytesIO(data))
        except Exception as exc:
            raise ValidationError(f"Could not read the PPTX file: {exc}") from exc

        pages: list[ParsedPage] = []
        for index, slide in enumerate(deck.slides, start=1):
            parts: list[str] = []
            title = None
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    parts.append(shape.text_frame.text.strip())
            if slide.shapes.title is not None and slide.shapes.title.text.strip():
                title = slide.shapes.title.text.strip()
            # Speaker notes carry the argument the slide only gestures at, and
            # are frequently the most answer-bearing text in a deck.
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append(f"Speaker notes: {notes}")
            text = normalize("\n\n".join(parts))
            if text:
                pages.append(ParsedPage(number=index, text=text, section=title))
        return ParsedDocument(pages=pages)


class TextParser:
    """Markdown and plain text.

    The heading tree drives the split directly, so Markdown produces the
    best-structured chunks of any format here.
    """

    source_type = SourceType.MARKDOWN
    _HEADING = re.compile(r"^(#{1,6})\s+(.*)$")

    def parse(self, data: bytes, *, filename: str | None = None) -> ParsedDocument:
        try:
            raw = data.decode("utf-8")
        except UnicodeDecodeError:
            raw = data.decode("utf-8", errors="replace")

        sections: list[tuple[str | None, list[str]]] = [(None, [])]
        for line in raw.split("\n"):
            match = self._HEADING.match(line)
            if match:
                sections.append((match.group(2).strip(), []))
            else:
                sections[-1][1].append(line)

        pages = [
            ParsedPage(number=i, text=normalize("\n".join(body)), section=heading)
            for i, (heading, body) in enumerate(sections, start=1)
            if normalize("\n".join(body))
        ]
        title = next((h for h, _ in sections if h), None)
        return ParsedDocument(pages=pages, title=title)


class HtmlParser:
    """Fetched web pages. The fetch itself is SSRF-guarded elsewhere (§18)."""

    source_type = SourceType.URL

    def parse(self, data: bytes, *, filename: str | None = None) -> ParsedDocument:
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(data, "html.parser")
        # Chrome, navigation and boilerplate are noise that dilutes retrieval:
        # every page in a site shares them, so they add no discriminating signal
        # and crowd out the content that does.
        for tag in soup(
            ["script", "style", "nav", "footer", "header", "aside", "noscript", "form"]
        ):
            tag.decompose()

        title = (soup.title.string or "").strip() if soup.title and soup.title.string else None
        body = soup.body or soup
        text = normalize(body.get_text(separator="\n"))
        pages = [ParsedPage(number=1, text=text)] if text else []
        return ParsedDocument(pages=pages, title=title)


# ── registry ─────────────────────────────────────────────────────────────
# Adding a format means adding a class above and one line here. The pipeline is
# never touched (Open/Closed, §6).

_BY_SOURCE_TYPE: dict[SourceType, object] = {
    SourceType.PDF: PdfParser(),
    SourceType.DOCX: DocxParser(),
    SourceType.PPTX: PptxParser(),
    SourceType.MARKDOWN: TextParser(),
    SourceType.TXT: TextParser(),
    SourceType.URL: HtmlParser(),
}

_BY_EXTENSION: dict[str, SourceType] = {
    ".pdf": SourceType.PDF,
    ".docx": SourceType.DOCX,
    ".pptx": SourceType.PPTX,
    ".md": SourceType.MARKDOWN,
    ".markdown": SourceType.MARKDOWN,
    ".txt": SourceType.TXT,
    ".text": SourceType.TXT,
    ".html": SourceType.URL,
    ".htm": SourceType.URL,
}

# Magic bytes, checked instead of trusting the extension (§18). A .pdf that is
# really a zip is either a mistake or an attack; both deserve rejection.
_MAGIC: tuple[tuple[bytes, SourceType], ...] = (
    (b"%PDF-", SourceType.PDF),
    (b"PK\x03\x04", SourceType.DOCX),  # OOXML is a zip; refined below
)


def source_type_for(filename: str | None, data: bytes) -> SourceType:
    """Determine the format from content first, filename second.

    OOXML formats share the zip magic, so the extension disambiguates DOCX from
    PPTX — but only after the content has confirmed it really is a zip.
    """
    head = data[:8]

    if head.startswith(b"%PDF-"):
        return SourceType.PDF

    extension = ""
    if filename and "." in filename:
        extension = "." + filename.rsplit(".", 1)[-1].lower()

    if head.startswith(b"PK\x03\x04"):
        if extension == ".pptx":
            return SourceType.PPTX
        if extension == ".docx":
            return SourceType.DOCX
        raise UnsupportedMediaError(
            "Zip-based file with an unrecognised extension; expected .docx or .pptx."
        )

    if extension in _BY_EXTENSION:
        declared = _BY_EXTENSION[extension]
        # A text-ish extension on binary content is a mislabelled file.
        if declared in {SourceType.PDF, SourceType.DOCX, SourceType.PPTX}:
            raise UnsupportedMediaError(
                f"File claims to be {extension} but its content does not match."
            )
        return declared

    # No extension and no magic: accept only if it decodes as text.
    try:
        data[:4096].decode("utf-8")
    except UnicodeDecodeError:
        raise UnsupportedMediaError(
            "Unsupported file type. Accepted: PDF, DOCX, PPTX, Markdown, text."
        ) from None
    return SourceType.TXT


def get_parser(source_type: SourceType):
    parser = _BY_SOURCE_TYPE.get(source_type)
    if parser is None:
        raise UnsupportedMediaError(f"No parser registered for {source_type.value}.")
    return parser
